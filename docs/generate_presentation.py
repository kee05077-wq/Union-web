"""
Union-web 프로젝트 소개 PPT 생성 스크립트.

실행:
    pip install python-pptx
    python generate_presentation.py

출력: Union-web_프로젝트소개.pptx (이 스크립트와 같은 폴더)
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ==========================================
# Union-web 브랜드 컬러 (static/css/main.css 기준)
# ==========================================
NAVY = RGBColor(0x12, 0x16, 0x2B)
NAVY_2 = RGBColor(0x1B, 0x21, 0x40)
CORAL = RGBColor(0xFF, 0x6B, 0x4A)
TEAL = RGBColor(0x2B, 0xB3, 0xA0)
SLATE = RGBColor(0x9C, 0xA3, 0xAF)
WHITE = RGBColor(0xFC, 0xFA, 0xF7)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUTPUT_PATH = Path(__file__).resolve().parent / "Union-web_프로젝트소개.pptx"


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, color=WHITE,
                 bold=False, align=PP_ALIGN.LEFT, font_name="맑은 고딕", anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_bullets(slide, left, top, width, height, items, size=18, color=WHITE,
                 font_name="맑은 고딕", line_spacing=1.3):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        level = 0
        text = item
        if isinstance(item, tuple):
            text, level = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.line_spacing = line_spacing
        run = p.add_run()
        prefix = "▸ " if level == 0 else "-  "
        run.text = f"{prefix}{text}"
        run.font.size = Pt(size - level * 2)
        run.font.color.rgb = color if level == 0 else SLATE
        run.font.name = font_name
    return box


def add_accent_bar(slide, color=TEAL):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_page_header(slide, kicker, title):
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(11.5), Inches(0.4),
                kicker, size=15, color=TEAL, bold=True)
    add_textbox(slide, Inches(0.7), Inches(0.72), Inches(11.9), Inches(0.8),
                title, size=30, color=WHITE, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.55), Inches(3.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = CORAL
    line.line.fill.background()


def new_slide(prs, bg=NAVY_2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_background(slide, bg)
    add_accent_bar(slide)
    return slide


def add_box(slide, left, top, width, height, text, fill=NAVY, line_color=TEAL,
            text_color=WHITE, size=14, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = "맑은 고딕"
    return shape


def add_arrow(slide, start_shape, end_shape, color=TEAL):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        start_shape.left + start_shape.width, start_shape.top + start_shape.height // 2,
        end_shape.left, end_shape.top + end_shape.height // 2,
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2.25)
    line = connector.line._get_or_add_ln()
    arrow = line.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    line.append(arrow)
    return connector


def add_file_table(slide, left, top, width, rows, col1_w_ratio=0.30, size=11, row_h=Inches(0.32)):
    """(파일명, 설명) 튜플 리스트를 표로 렌더링한다."""
    n = len(rows)
    height = row_h * n
    table_shape = slide.shapes.add_table(n, 2, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = int(width * col1_w_ratio)
    table.columns[1].width = width - table.columns[0].width

    for r, (name, desc) in enumerate(rows):
        table.rows[r].height = row_h

        cell0 = table.cell(r, 0)
        cell0.fill.solid()
        cell0.fill.fore_color.rgb = NAVY
        cell0.margin_top = cell0.margin_bottom = Pt(2)
        cell0.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = cell0.text_frame.paragraphs[0]
        run0 = p0.add_run()
        run0.text = name
        run0.font.size = Pt(size)
        run0.font.bold = True
        run0.font.color.rgb = TEAL
        run0.font.name = "Consolas"

        cell1 = table.cell(r, 1)
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = NAVY_2
        cell1.margin_top = cell1.margin_bottom = Pt(2)
        cell1.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = cell1.text_frame.paragraphs[0]
        run1 = p1.add_run()
        run1.text = desc
        run1.font.size = Pt(size)
        run1.font.color.rgb = WHITE
        run1.font.name = "맑은 고딕"

    # 표 기본 스타일(줄무늬 등) 제거
    tbl = table_shape.table._tbl
    tblPr = tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')

    return table_shape


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---------- 1. 타이틀 ----------
    slide = new_slide(prs, bg=NAVY)
    add_textbox(slide, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
                "Union-web", size=54, color=WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(3.55), Inches(11.3), Inches(0.8),
                "AI 수어 인식 기반 실시간 화상회의 플랫폼", size=24, color=TEAL, bold=True)
    add_textbox(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6),
                "청각장애인과 비장애인이 자막 없이도 자연스럽게 소통할 수 있는 화상회의 서비스",
                size=16, color=SLATE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.45), Inches(2.4), Pt(4))
    line.fill.solid(); line.fill.fore_color.rgb = CORAL; line.line.fill.background()

    # ---------- 2. 프로젝트 개요 ----------
    slide = new_slide(prs)
    add_page_header(slide, "PROJECT OVERVIEW", "프로젝트 개요")
    add_bullets(slide, Inches(0.9), Inches(2.0), Inches(11.3), Inches(4.5), [
        "문제의식: 화상회의에서 수어 사용자는 실시간으로 의사를 전달하기 어렵고, 별도 통역 인력 없이는 소통이 제한됨",
        "목표: 카메라로 촬영한 수어(지문자)를 AI가 실시간으로 인식해 자동으로 한글 자막을 생성/전송",
        "핵심 가치",
        ("실시간성 - 화상회의 도중 지연 없이 자막 생성", 1),
        ("접근성 - 별도 장비 없이 웹캠 + 브라우저만으로 사용 가능", 1),
        ("확장성 - 지문자 인식에서 시작해 단어/문장 단위 인식으로 확장 예정", 1),
    ], size=19)

    # ---------- 3. 시스템 아키텍처 ----------
    slide = new_slide(prs)
    add_page_header(slide, "SYSTEM ARCHITECTURE", "전체 시스템 구성도")
    boxes_y = Inches(2.3)
    labels = ["브라우저\n(WebRTC 화상통화)", "Flask-SocketIO\n서버", "MediaPipe\n손 랜드마크 검출", "LSTM 모델\n(자모 예측)", "자모→음절 조합\n실시간 자막"]
    bw, bh, gap = Inches(2.15), Inches(1.3), Inches(0.28)
    x = Inches(0.55)
    prev = None
    for label in labels:
        box = add_box(slide, x, boxes_y, bw, bh, label, size=13)
        if prev is not None:
            add_arrow(slide, prev, box)
        prev = box
        x = Emu(x + bw + gap)
    add_bullets(slide, Inches(0.9), Inches(4.3), Inches(11.3), Inches(2.6), [
        "각 참가자의 카메라 프레임을 서버로 전송 -> 서버에서 MediaPipe로 손 랜드마크 추출",
        "10프레임 시퀀스를 LSTM 모델에 입력해 자모(지문자) 예측",
        "예측된 자모를 음절로 조합해 실시간 자막(sign_progress)과 완성 문장(sign_result)으로 송출",
        "무거운 연산(MediaPipe/TensorFlow)은 별도 스레드(eventlet tpool)에서 처리해 여러 참가자가 동시에 접속해도 서버가 멈추지 않도록 구성",
    ], size=16)

    # ---------- 4. 핵심 기능 ----------
    slide = new_slide(prs)
    add_page_header(slide, "KEY FEATURES", "핵심 기능")
    feature_items = [
        ("회원 관리", "회원가입 / 로그인 / 아이디·비밀번호 찾기"),
        ("실시간 화상회의", "WebRTC 기반 1:1 화상통화, 음성/영상 On-Off"),
        ("실시간 채팅", "텍스트 및 파일 공유 채팅"),
        ("실시간 수어 자막", "손 모양을 인식해 자동으로 한글 자막 생성 및 상대방에게 전송"),
    ]
    y = Inches(2.1)
    for title, desc in feature_items:
        add_box(slide, Inches(0.9), y, Inches(2.6), Inches(0.95), title, size=15)
        add_textbox(slide, Inches(3.75), y, Inches(8.5), Inches(0.95), desc, size=15, color=WHITE,
                    anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(y + Inches(1.15))

    # ---------- 5. 수어 인식 파이프라인 상세 ----------
    slide = new_slide(prs)
    add_page_header(slide, "RECOGNITION PIPELINE", "수어 인식 파이프라인 상세")
    steps = [
        "카메라 프레임 캡처 (약 35fps, 브라우저 -> 서버 전송)",
        "MediaPipe로 오른손 21개 랜드마크 추출",
        "관절 벡터 20개 + 관절 각도 15개 = 55차원 특징 벡터 계산",
        "10프레임 시퀀스로 누적",
        "LSTM 모델로 자모(31개 클래스 중 1개) 예측 + 신뢰도 확인",
        "자음/모음 전환 감지 시 직전 음운을 확정 -> 음절 조합 (join_jamos)",
        "완성된 단어를 문장으로 묶어 실시간 자막으로 상대방에게 전송",
    ]
    y = Inches(1.95)
    for i, step in enumerate(steps, start=1):
        add_textbox(slide, Inches(0.9), y, Inches(0.5), Inches(0.5), str(i), size=18, color=CORAL, bold=True)
        add_textbox(slide, Inches(1.4), y, Inches(10.8), Inches(0.5), step, size=16, color=WHITE)
        y = Emu(y + Inches(0.68))

    # ---------- 6. 모델 선택 배경 (YOLO vs LSTM) ----------
    slide = new_slide(prs, bg=NAVY)
    add_page_header(slide, "MODEL DESIGN DECISION", "왜 YOLO가 아닌 LSTM 시퀀스 모델인가")
    add_bullets(slide, Inches(0.9), Inches(2.0), Inches(11.3), Inches(4.5), [
        "초기에는 YOLO 기반 전이학습(객체 탐지)으로 손모양을 인식하는 방식을 검토함",
        "그러나 YOLO는 정지된 한 프레임 안에서 바운딩박스로 손 모양을 탐지하는 방식이라, 실시간 영상에서 이어지는 손동작의 시간적 흐름(시퀀스)을 반영하지 못함",
        "화상회의처럼 끊김 없는 실시간 처리가 핵심인 서비스에는, 프레임 간 흐름을 학습하는 LSTM 기반 시퀀스 모델이 더 적합하다고 판단해 채택",
    ], size=20)

    # ---------- 7. 현재 상태 ----------
    slide = new_slide(prs)
    add_page_header(slide, "CURRENT STATUS", "현재 상태")
    add_bullets(slide, Inches(0.9), Inches(2.0), Inches(11.3), Inches(4.5), [
        "현재 배포된 인식 모델은 파이프라인 검증을 위해 공개된 모델을 가져와 연동한 것으로, 자체 데이터로 학습한 모델이 아님",
        "지문자(자음/모음) 31종만 인식 가능 - 실제 대화에 필요한 단어/문장 단위 수어는 아직 지원하지 않음",
        "따라서 다음 단계로 자체 데이터셋을 구축하고 모델을 직접 학습시키는 과정이 필요함",
    ], size=20)

    # ---------- 8. 기술적 한계 ----------
    slide = new_slide(prs, bg=NAVY)
    add_page_header(slide, "TECHNICAL LIMITATIONS", "기술적 한계")
    add_bullets(slide, Inches(0.9), Inches(1.95), Inches(11.3), Inches(4.9), [
        "시퀀스 대기 지연: LSTM 모델은 10프레임이 모여야 예측하므로, 사용자가 손모양을 유지한 채 잠시 기다려야 함",
        ("이론상 35fps면 약 0.3초면 충분하지만, 실제로는 프레임 캡처 후 서버로 전송하는 통신 시간이 더해져 체감 지연이 발생", 1),
        "자모 조합 알고리즘의 한계: 실시간 영상 특성상 손이 다음 동작으로 이동하는 전환 구간에서, 사용자가 의도하지 않은 자모가 순간적으로 인식되는 경우가 있음",
        "팀 구성의 한계: 팀 전원이 모델 학습 경험이 전무한 1학년으로 구성되어, 학습에 필요한 하이퍼파라미터(신뢰도 임계값, 시퀀스 길이, epoch 수 등)를 충분한 근거 없이 시행착오로 정해야 했음",
    ], size=18)

    # ---------- 9. 향후계획 1: 자체 모델 학습 ----------
    slide = new_slide(prs)
    add_page_header(slide, "FUTURE PLAN 1", "자체 데이터셋 구축 및 모델 학습")
    add_textbox(slide, Inches(0.9), Inches(1.95), Inches(11), Inches(0.5),
                "수어 방식에 따라 서로 다른 특징과 구조를 갖는 2가지 모델을 별도로 학습한다", size=16, color=SLATE)

    col_w = Inches(5.6)
    add_box(slide, Inches(0.9), Inches(2.6), col_w, Inches(0.7), "① 음운(지문자) 인식 모델", fill=NAVY, size=16)
    add_bullets(slide, Inches(0.9), Inches(3.45), col_w, Inches(3.2), [
        "대상: 자음/모음 낱자 31종",
        "입력: 오른손 21개 랜드마크 (55차원 특징, 10프레임 시퀀스)",
        "구조: LSTM(64) -> Dropout -> Dense(32) -> Dense(softmax)",
    ], size=15)

    add_box(slide, Inches(6.85), Inches(2.6), col_w, Inches(0.7), "② 단어 인식 모델", fill=NAVY, size=16)
    add_bullets(slide, Inches(6.85), Inches(3.45), col_w, Inches(3.2), [
        "대상: 나는 / 맛있다 / 고양이 / 식사 등 단어 단위 수어",
        "입력: 양손 21점x2 + 어깨·팔꿈치·손목 6점 (150차원 특징, 30프레임 시퀀스)",
        "구조: LSTM(64,seq) -> LSTM(64) -> Dropout -> Dense(32) -> Dense(softmax)",
    ], size=15)

    # ---------- 10. 데이터 수집 방법 ----------
    slide = new_slide(prs)
    add_page_header(slide, "FUTURE PLAN 1 (cont.)", "데이터 수집 방법 2가지")
    add_box(slide, Inches(0.9), Inches(2.1), Inches(5.6), Inches(0.7), "방법 1. 직접 촬영", size=16)
    add_bullets(slide, Inches(0.9), Inches(2.95), Inches(5.6), Inches(3.2), [
        "웹캠으로 클래스별 영상을 직접 녹화",
        "MediaPipe로 랜드마크 추출 후 시퀀스(npy)로 변환",
        "소규모 클래스 확장, 특정 상황 데이터 보강에 적합",
    ], size=15)
    add_box(slide, Inches(6.85), Inches(2.1), Inches(5.6), Inches(0.7), "방법 2. AIHub 공개 데이터", size=16)
    add_bullets(slide, Inches(6.85), Inches(2.95), Inches(5.6), Inches(3.2), [
        "AIHub “수어 영상” 데이터셋 활용",
        "지화 1,000종 + 수어 단어 3,000종 + 수어 문장 2,000종, 형태소 단위 라벨링 제공",
        "대규모 데이터 확보로 모델 일반화 성능 향상에 활용",
    ], size=15)
    add_textbox(slide, Inches(0.9), Inches(6.35), Inches(11.3), Inches(0.5),
                "출처: aihub.or.kr - 수어 영상 데이터셋 (dataSetSn=103)", size=13, color=SLATE)

    # ---------- 11. 향후계획 2: 텍스트/음성 -> 수어 애니메이션 ----------
    slide = new_slide(prs, bg=NAVY)
    add_page_header(slide, "FUTURE PLAN 2", "텍스트·음성 → 수어 애니메이션 변환")
    add_bullets(slide, Inches(0.9), Inches(2.1), Inches(11.3), Inches(4.3), [
        "현재는 수어 -> 텍스트(자막) 방향만 지원",
        "반대 방향으로, 입력된 텍스트 또는 음성을 3D 아바타의 수어 동작으로 변환하는 기능을 계획",
        "청각장애인이 상대방의 말(텍스트/음성)을 수어 애니메이션으로 시각적으로 확인 가능",
        "-> 화상회의 내 양방향(수어 ↔ 텍스트/음성) 실시간 번역 완성",
    ], size=19)

    # ---------- 12. 향후계획 3: 음성 -> 전사문 ----------
    slide = new_slide(prs)
    add_page_header(slide, "FUTURE PLAN 3", "음성 → 전사문(STT) 변환")
    add_bullets(slide, Inches(0.9), Inches(2.1), Inches(11.3), Inches(4.3), [
        "상대방(비장애인)의 음성을 실시간 음성인식(STT)으로 텍스트化",
        "생성된 전사문은 자막으로 표시하거나, 향후계획 2(수어 애니메이션 변환)의 입력으로 사용",
        "결과적으로 '음성 -> 전사문 -> 수어 애니메이션'으로 이어지는 파이프라인의 첫 단계",
    ], size=19)

    # ---------- 13. 향후계획 4: 자연스러운 문장 변환 ----------
    slide = new_slide(prs)
    add_page_header(slide, "FUTURE PLAN 4", "단어 나열 → 자연스러운 문장 변환")
    add_bullets(slide, Inches(0.9), Inches(2.1), Inches(11.3), Inches(4.3), [
        "수어는 조사·어미 없이 핵심 단어만 순서대로 나열하는 방식으로 표현되는 경우가 많음",
        "예) “나 / 학교 / 가다” 처럼 인식된 단어열은 그 자체로는 자연스러운 한국어 문장이 아님",
        "인식된 단어열을 자연어 생성(언어 모델)을 통해 조사·어미가 포함된 자연스러운 문장으로 재구성하는 장치가 필요",
        "예) “나 학교 가다” -> “저는 학교에 갑니다”",
    ], size=19)

    # ---------- 14. 향후계획 5: 수어 -> 음성 변환 (예고) ----------
    slide = new_slide(prs, bg=NAVY)
    add_page_header(slide, "FUTURE PLAN 5 (예고)", "수어 → 음성 변환 장치")
    add_bullets(slide, Inches(0.9), Inches(2.1), Inches(11.3), Inches(4.3), [
        "인식된 수어 자막(문장)을 음성(TTS)으로 변환해 상대방에게 들려주는 기능도 함께 검토 중",
        "화상회의뿐 아니라 오프라인 대면 상황에서도 활용 가능한 방향으로 확장 고려",
        "구체적인 설계는 앞선 4가지 향후계획 이후 순차적으로 진행 예정",
    ], size=19)

    # ---------- 15. 로드맵 요약 ----------
    slide = new_slide(prs)
    add_page_header(slide, "ROADMAP", "전체 로드맵 요약")
    roadmap = [
        ("현재", "지문자 인식 (외부 모델) + 화상회의/채팅"),
        ("1단계", "자체 데이터셋 구축 및 음운/단어 모델 학습"),
        ("2단계", "음성 → 전사문(STT) 및 자연스러운 문장 변환"),
        ("3단계", "텍스트·음성 → 수어 애니메이션 변환 (양방향 번역 완성)"),
        ("4단계", "수어 → 음성 변환 장치 검토"),
    ]
    y = Inches(2.0)
    for stage, desc in roadmap:
        add_box(slide, Inches(0.9), y, Inches(1.7), Inches(0.85), stage, fill=CORAL, line_color=CORAL, size=16)
        add_textbox(slide, Inches(2.85), y, Inches(9.4), Inches(0.85), desc, size=16, color=WHITE,
                    anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(y + Inches(1.0))

    # ---------- 16. 부록: Union-web 폴더 구조 ----------
    slide = new_slide(prs)
    add_page_header(slide, "APPENDIX", "부록: Union-web 폴더 구조")
    rows = [
        ("app.py", "Flask+Socket.IO 서버 - 라우팅, 회원 DB, 수어 인식(MediaPipe+LSTM), 자모 조합 로직"),
        ("requirements.txt", "의존 패키지 목록 (Flask, mediapipe, tensorflow 등 버전 고정)"),
        ("unicode.py", "자모 결합/분리 유틸 (join_jamos 등)"),
        ("gesture_classifier.h5", "배포된 지문자 인식 LSTM 모델 (외부에서 가져온 모델)"),
        ("modules/hand_module.py", "MediaPipe HandLandmarker 래퍼 (오른손 랜드마크 검출)"),
        ("modules/utils.py", "손 랜드마크 -> 55차원 특징 벡터 변환(Vector_Normalization)"),
        ("models/hand_landmarker.task", "MediaPipe 손 검출 모델 번들 파일"),
        ("database/users.xlsx", "회원 정보 저장 엑셀 DB"),
        ("templates/index.html 등", "랜딩/회원가입/계정찾기/소개 페이지"),
        ("templates/chatentry.html", "화상회의 입장 전 대기실 (카메라 미리보기)"),
        ("templates/room.html", "화상회의방 - WebRTC 화상통화 + 실시간 자막"),
        ("static/js/room_ui.js", "화상회의방 클라이언트 - WebRTC, 프레임 캡처/전송, 자막 표시"),
        ("static/js/chatentry.js", "입장 전 카메라/마이크 미리보기 로직"),
        ("static/js/auth.js 등", "로그인/회원가입/계정찾기 클라이언트 로직"),
    ]
    add_file_table(slide, Inches(0.7), Inches(1.9), Inches(11.9), rows, size=12, row_h=Inches(0.365))

    # ---------- 17. 부록: KSL-Model-Training 폴더 구조 ----------
    slide = new_slide(prs)
    add_page_header(slide, "APPENDIX", "부록: KSL-Model-Training 폴더 구조")
    rows2 = [
        ("README.md", "프로젝트 사용법 전체 안내"),
        ("requirements.txt", "의존 패키지 (mediapipe, tensorflow, scikit-learn 등)"),
        ("download_mediapipe_models.py", "MediaPipe HolisticLandmarker 모델 번들 다운로드"),
        ("modules/holistic_module.py", "MediaPipe HolisticLandmarker 래퍼 (양손+pose 검출)"),
        ("modules/utils.py", "음운용(55차원)/단어용(150차원) 특징 추출 함수"),
        ("phoneme_model/config.py", "음운 모델 설정 (31개 클래스, 시퀀스 길이 10 등)"),
        ("phoneme_model/record_dataset.py", "[방법1] 웹캠으로 지문자 영상 직접 녹화"),
        ("phoneme_model/build_dataset_from_videos.py", "녹화 영상 -> 학습용 시퀀스(npy) 변환"),
        ("phoneme_model/build_dataset_from_aihub.py", "[방법2] AIHub 공개 데이터 -> 시퀀스 변환"),
        ("phoneme_model/train.py", "음운 인식 LSTM 모델 학습"),
        ("phoneme_model/test_webcam.py", "웹캠으로 음운 모델 실시간 테스트"),
        ("word_model/*.py", "위와 동일 구성 - 단어 인식 모델용 (150차원, 30프레임 시퀀스)"),
    ]
    add_file_table(slide, Inches(0.7), Inches(1.9), Inches(11.9), rows2, size=13, row_h=Inches(0.42))

    # ---------- 18. 마무리 ----------
    slide = new_slide(prs, bg=NAVY)
    add_textbox(slide, Inches(1), Inches(3.1), Inches(11.3), Inches(1.2),
                "감사합니다", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(4.1), Inches(11.3), Inches(0.6),
                "Union-web Team", size=18, color=TEAL, align=PP_ALIGN.CENTER)

    prs.save(OUTPUT_PATH)
    print(f"저장됨: {OUTPUT_PATH}")


if __name__ == "__main__":
    build()
